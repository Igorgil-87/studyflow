"""
rag/document_extractor.py — extrai texto (em Markdown) de arquivos de
material de apoio pra indexar no RAG. Usado SÓ pelo módulo Curso — os
outros módulos (Youtuber, Criador, etc.) continuam sem essa dependência.

Motor principal: MarkItDown (Microsoft) — cobre PDF, PPTX, DOCX, XLSX,
XLS, CSV, JSON, XML, HTML, ZIP (itera o conteúdo), EPUB, imagens (EXIF +
OCR se markitdown-ocr estiver instalado), áudio (transcrição, se as libs
de audio-transcription estiverem instaladas). Testado de verdade com
PDF/PPTX/DOCX/XLSX antes de entrar aqui — sem conflito de dependência
com o resto do projeto (openai==1.51.0 preservado).

Fallback: se o markitdown não estiver instalado, cai pros extratores
manuais antigos (pypdf/python-pptx/python-docx) — só PDF/PPTX/DOCX, mas
continua funcionando sem quebrar nada que já existia.

Cada função aceita bytes (conteúdo do arquivo já lido) e devolve o texto
extraído em Markdown. Import tardio das libs — se não estiver instalada,
o erro só aparece quando alguém tentar usar aquele formato, não trava o
resto do app.
"""

from __future__ import annotations

import io

# Formatos que o MarkItDown cobre quando instalado (pip install
# "markitdown[pdf,pptx,docx,xlsx,xls]" — extras escolhidos a dedo, sem
# puxar Azure Document Intelligence/Content Understanding, que ninguém
# usa aqui). Imagem e áudio entram na lista mas dependem de extras que
# não instalamos por padrão (OCR e transcrição) — MarkItDown ainda
# extrai metadados EXIF de imagem sem eles, só não faz OCR/transcrição.
_MARKITDOWN_EXTENSIONS = {
    ".pdf", ".pptx", ".docx", ".xlsx", ".xls",
    ".csv", ".json", ".xml", ".html", ".htm",
    ".zip", ".epub",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp",
    ".mp3", ".wav", ".m4a",
}
# Formatos do fallback manual (sempre disponíveis — 0 dependência extra
# além do que já estava no requirements.txt antes do MarkItDown existir).
_FALLBACK_EXTENSIONS = {".pdf", ".pptx", ".docx"}


# Texto puro — decodifica direto, sem depender de MarkItDown nem de
# nenhuma lib de conversão. Sempre suportado, com ou sem MarkItDown
# instalado (por isso fica fora de _MARKITDOWN_EXTENSIONS/_FALLBACK_EXTENSIONS).
_PLAINTEXT_EXTENSIONS = {".txt", ".md", ".markdown"}


class DocumentExtractionError(RuntimeError):
    """Erro ao extrair texto de um documento."""


def _markitdown_available() -> bool:
    try:
        import markitdown  # noqa: F401
        return True
    except ImportError:
        return False


def _supported_extensions() -> set[str]:
    base = _MARKITDOWN_EXTENSIONS if _markitdown_available() else _FALLBACK_EXTENSIONS
    return base | _PLAINTEXT_EXTENSIONS


# Mantido como atributo de módulo (não função) porque a rota do Curso
# importa isso direto — recalculado a cada import, então reflete se o
# markitdown foi instalado ou não nesse ambiente.
SUPPORTED_EXTENSIONS = _supported_extensions()


def extract_text(filename: str, content: bytes) -> str:
    """Detecta o formato pela extensão do arquivo e extrai o texto (em
    Markdown, quando via MarkItDown). Levanta DocumentExtractionError se
    o formato não for suportado ou a extração falhar."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in _PLAINTEXT_EXTENSIONS:
        return _extract_plaintext(content)

    if _markitdown_available() and ext in _MARKITDOWN_EXTENSIONS:
        try:
            return _extract_with_markitdown(filename, content, ext)
        except DocumentExtractionError:
            # MarkItDown está instalado mas pode faltar o EXTRA específico
            # desse formato (ex: markitdown sem "[docx]") — se o fallback
            # manual cobrir esse formato, tenta antes de desistir.
            if ext in _FALLBACK_EXTENSIONS:
                try:
                    return _extract_fallback(ext, content)
                except DocumentExtractionError:
                    pass  # fallback também falhou -> segue e relança o erro original
            raise

    # sem MarkItDown instalado (ou formato fora da lista dele) -> tenta o
    # fallback manual, que só cobre PDF/PPTX/DOCX
    if ext in _FALLBACK_EXTENSIONS:
        return _extract_fallback(ext, content)

    if ext in _MARKITDOWN_EXTENSIONS:
        raise DocumentExtractionError(
            f"Formato '{ext}' precisa do MarkItDown, que não está instalado "
            f'(pip install "markitdown[pdf,pptx,docx,xlsx,xls]").'
        )
    raise DocumentExtractionError(
        f"Formato '{ext}' não suportado. Formatos aceitos: "
        f"{', '.join(sorted(_supported_extensions()))}"
    )


def _extract_plaintext(content: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = content.decode(encoding).strip()
            break
        except UnicodeDecodeError:
            continue
    else:
        raise DocumentExtractionError(
            "Não consegui decodificar esse arquivo de texto (encoding não reconhecido)."
        )
    if not text:
        raise DocumentExtractionError("Arquivo de texto vazio.")
    return text


def _extract_fallback(ext: str, content: bytes) -> str:
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".pptx":
        return _extract_pptx(content)
    if ext == ".docx":
        return _extract_docx(content)
    raise DocumentExtractionError(f"Sem extrator de fallback pra '{ext}'.")


def _extract_with_markitdown(filename: str, content: bytes, ext: str) -> str:
    try:
        from markitdown import MarkItDown
    except ImportError as exc:
        raise DocumentExtractionError(
            'markitdown não instalado (pip install "markitdown[pdf,pptx,docx,xlsx,xls]")'
        ) from exc

    try:
        md = MarkItDown(enable_plugins=False)
        result = md.convert_stream(io.BytesIO(content), file_extension=ext)
    except Exception as exc:
        raise DocumentExtractionError(f"Falha ao converter '{filename}': {exc}") from exc

    text = (result.markdown or "").strip()
    if not text:
        raise DocumentExtractionError(
            f"Não consegui extrair texto de '{filename}' — arquivo vazio, "
            "corrompido, ou (se for imagem/PDF escaneado) precisa do "
            'extra de OCR: pip install markitdown-ocr'
        )
    return text


# ── Fallback manual (PDF/PPTX/DOCX) — usado só se o markitdown não
# estiver instalado nesse ambiente. Mantido por compatibilidade; o
# caminho principal agora é _extract_with_markitdown() acima. ──────────

def _extract_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentExtractionError("pypdf não instalado (pip install pypdf)") from exc

    try:
        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise DocumentExtractionError(f"Falha ao ler o PDF: {exc}") from exc

    text = "\n\n".join(p.strip() for p in pages if p.strip())
    if not text:
        raise DocumentExtractionError(
            "Não consegui extrair texto desse PDF — provavelmente é um PDF "
            "escaneado (imagem), sem camada de texto. Pra isso funcionar, "
            'instala o MarkItDown com OCR: pip install "markitdown[pdf]" markitdown-ocr'
        )
    return text


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentExtractionError("python-pptx não instalado (pip install python-pptx)") from exc

    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise DocumentExtractionError(f"Falha ao ler o PPTX: {exc}") from exc

    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        if parts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))

    text = "\n\n".join(slides_text)
    if not text:
        raise DocumentExtractionError("Não encontrei texto nesse PPTX (slides vazios ou só imagens).")
    return text


def _extract_docx(content: bytes) -> str:
    try:
        import docx
    except ImportError as exc:
        raise DocumentExtractionError("python-docx não instalado (pip install python-docx)") from exc

    try:
        doc = docx.Document(io.BytesIO(content))
    except Exception as exc:
        raise DocumentExtractionError(f"Falha ao ler o DOCX: {exc}") from exc

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    if not text:
        raise DocumentExtractionError("Não encontrei texto nesse DOCX (documento vazio).")
    return text
